# !/usr/bin/env python
# -*-coding:utf-8 -*-

"""
# File       : file_io.py
# Time       ：2024/10/8 14:06
# Author     ：jianbang
# version    ：python 3.10
# company    : IFLYTEK Co.,Ltd.
# emil       : whdx072018@foxmail.com
# Description：
"""
import os
import zipfile
from PIL import ImageDraw, ImageFont
from docx import Document
import pdf2image
from PyPDF2 import PdfReader
from pptx import Presentation
from PIL import Image
import subprocess
from natsort import natsorted
import pdfplumber
import platform
import pandas as pd
from openpyxl import load_workbook
import xlrd



class FileIOApi(object):
    def write_text(self,content,txt_path):
        """
        :param content:
        :param txt_path:
        :return:
        """
        with open(txt_path,"w",encoding="utf-8") as txt_file:
            txt_file.write(content)
        print(f"write data into: {txt_path}")

    def read_pdf(self,file_path):
        """
        :param file_path:
        :return:
        """
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            content = ""
            for page in reader.pages:
                content += page.extract_text()
            return content

    def read_docx(self, file_path):
        """
        :param file_path:
        :return:
        """
        doc = Document(file_path)
        content = []
        for para in doc.paragraphs:
            content.append(para.text)
        return '\n'.join(content)

    def read_pptx(self,file_path):
        """
        :param file_path:
        :return:
        """
        prs = Presentation(file_path)
        content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        return '\n'.join(content)

    def read_xlsx(self, file_path):
        """
        :param file_path:
        :return:
        """
        return pd.read_excel(file_path)

    def docx_to_pdf(self, docx_file_path, pdf_file_path):
        """
        :param docx_file_path:
        :param pdf_file_path:
        :return:
        """
        current_os = platform.system()

        if current_os == 'Darwin':
            libreoffice_path = '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        elif current_os == 'Linux':
            libreoffice_path = 'libreoffice'
        else:
            raise OSError("当前系统不支持。请确保系统为 macOS 或 Linux。")

        try:
            subprocess.run([libreoffice_path, '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(pdf_file_path),
                            docx_file_path])
            print(f"Converted {docx_file_path} to PDF successfully.")
        except subprocess.CalledProcessError as e:
            print(f"Error during conversion: {e}")

    def pptx_to_pdf(self,pptx_file_path, pdf_file_path):
        """
        :param pptx_file_path:
        :param pdf_file_path:
        :return:
        """
        current_os = platform.system()

        if current_os == 'Darwin':
            libreoffice_path = '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        elif current_os == 'Linux':
            libreoffice_path = 'libreoffice'
        else:
            raise OSError("当前系统不支持。请确保系统为 macOS 或 Linux。")

        try:
            subprocess.run([libreoffice_path, '--headless', '--convert-to', 'pdf', pptx_file_path, '--outdir',
                            os.path.dirname(pdf_file_path)], check=True)
            print(f"Converted {pptx_file_path} to PDF successfully.")
        except subprocess.CalledProcessError as e:
            print(f"Error during conversion: {e}")

    def pdf_all_to_image(self, pdf_file_path, img_dir, dpi=200):
        """
        :param pdf_file_path:
        :param output_folder:
        :param dpi:
        :return:
        """
        print("正在转化为image...")
        images = pdf2image.convert_from_path(pdf_file_path, dpi=dpi)
        os.makedirs(img_dir, exist_ok=True)
        for i, image in enumerate(images):
            image.save(os.path.join(img_dir, f'{i}.png'))

    def docx_all_to_image(self, docx_file_path, docx_pdf_dir, output_folder, dpi=200,transform_pdf_before=False,pdf_file_path=None):
        """
        :param docx_file_path:
        :param docx_pdf_dir:
        :param output_folder:
        :param dpi:
        :param transform_pdf_before:
        :param pdf_file_path:
        :return:
        """
        docx_file_name = os.path.basename(docx_file_path)
        pdf_file_name = docx_file_name.replace("docx", "pdf")
        filename = docx_file_name.replace(".docx", "").strip()
        img_dir = os.path.join(output_folder, filename)
        if not transform_pdf_before and pdf_file_path is None:
            pdf_file_path = os.path.join(docx_pdf_dir, pdf_file_name)
            self.docx_to_pdf(docx_file_path, pdf_file_path)

        self.pdf_all_to_image(pdf_file_path, img_dir, dpi)


    def pptx_all_to_image(self, pptx_file_path, pptx_pdf_dir, output_folder, dpi=200):
        """
        :param pptx_file_path:
        :param pptx_pdf_dir:
        :param output_folder:
        :param dpi:
        :return:
        """
        pptx_file_name = os.path.basename(pptx_file_path)
        pdf_file_name = pptx_file_name.replace("pptx","pdf")
        filename =pptx_file_name.replace(".pptx","").strip()
        img_dir = os.path.join(output_folder, filename)

        pdf_file_path = os.path.join(pptx_pdf_dir, pdf_file_name)
        self.pptx_to_pdf(pptx_file_path,pdf_file_path)
        self.pdf_all_to_image(pdf_file_path,img_dir,dpi)

    def pdf_page_to_image(self, pdf_file_path, output_folder, image_lst, dpi=200):
        """
        :param pdf_file_path:
        :param output_folder:
        :param image_lst:
        :param dpi:
        :return:
        """
        images = pdf2image.convert_from_path(pdf_file_path, dpi=dpi)
        file_name_with_extension = os.path.basename(pdf_file_path)
        file_name, _ = os.path.splitext(file_name_with_extension)
        for i, image in enumerate(images):
            if i not in image_lst:
                continue
            image.save(os.path.join(output_folder, f'{i}.png'))

    def docx_page_to_image(self, docx_file_path, output_folder, pages_lst, font_path=None, font_size=14,
                           image_width=1024, image_height=2048):
        """
        :param docx_file_path:
        :param output_folder:
        :param pages_lst:
        :param font_path:
        :param font_size:
        :param image_width:
        :param image_height:
        :return:
        """
        doc = Document(docx_file_path)
        file_name_with_extension = os.path.basename(docx_file_path)
        file_name, _ = os.path.splitext(file_name_with_extension)
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        img_dir = os.path.join(output_folder, file_name)
        os.makedirs(img_dir, exist_ok=True)

        for page_num, para in enumerate(doc.paragraphs, start=1):
            if page_num not in pages_lst:
                continue
            image = Image.new('RGB', (image_width, image_height), color=(255, 255, 255))
            draw = ImageDraw.Draw(image)

            if font_path:
                font = ImageFont.truetype(font_path, font_size)
            else:
                font = ImageFont.load_default()

            text_position = [10, 10]
            draw.text(text_position, para.text, font=font, fill=(0, 0, 0))

            image_file_name = f"page_{page_num}.png"
            image.save(os.path.join(img_dir, image_file_name))

    def pptx_page_to_image(self, pptx_file_path, pdf_file_path, output_folder,image_lst,dpi=200):
        """
        :param pptx_file_path:
        :param pdf_file_path:
        :param output_folder:
        :param image_lst:
        :param dpi:
        :return:
        """
        self.pptx_to_pdf(pptx_file_path, pdf_file_path)
        self.pdf_page_to_image(pdf_file_path, output_folder, image_lst,dpi=dpi)


    def read_pdf_pages(self, file_path, img_pages_lst,output_data_dir):
        """
        :param file_path:
        :param img_pages_lst:
        :param output_data_dir:
        :return:
        """

        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)

            pages_to_read = [i for i in range(total_pages) if i not in img_pages_lst]

            for page_num in pages_to_read:
                if page_num < total_pages:
                    page = pdf.pages[page_num]
                    page_text = page.extract_text()

                    txt_path=os.path.join(output_data_dir,str(page_num)+".txt")
                    if page_text:
                        self.write_text(page_text,txt_path)

    def read_docx_pages(self, file_path, img_pages_lst):
        """
        :param file_path:
        :param img_pages_lst:
        :return:
        """
        doc = Document(file_path)
        content = []

        for i, para in enumerate(doc.paragraphs):
            if i + 1 in img_pages_lst:
                continue
            content.append(para.text)
        return '\n'.join(content)

    def read_pptx_pages(self,file_path, img_pages_lst,output_data_dir):
        """
        :param file_path:
        :param img_pages_lst:
        :param output_data_dir:
        :return:
        """
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            if i in img_pages_lst:
                continue

            content=""
            txt_path = os.path.join(output_data_dir, str(i) + ".txt")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content+="\n"+shape.text
            self.write_text(content, txt_path)


    def extract_images_tables_from_pdf(self, pdf_file):
        """
        :param pdf_file: The path to the PDF file
        :return: A tuple containing lists of extracted images and tables from the PDF, and the total number of pages
        """
        images_in_pdf, tables_in_pdf = [], []

        with pdfplumber.open(pdf_file) as pdf:
            total_pages = len(pdf.pages)

            for page_num, page in enumerate(pdf.pages):
                header_height = page.height / 20

                width, height = page.width, page.height
                non_header_bbox = (0, header_height, width, height)

                # Extract content from the non-header region
                cropped_page = page.within_bbox(non_header_bbox)

                # Extract images from the non-header region
                images = cropped_page.images
                if images:
                    for img in images:
                        images_in_pdf.append({
                            'page_num': page_num + 1,
                            'image_data': img,
                        })

                # Extract tables from the non-header region
                tables = cropped_page.extract_tables()
                if tables:
                    for table in tables:
                        tables_in_pdf.append({
                            'page_num': page_num + 1,
                            'table_data': table,
                        })

        return images_in_pdf, tables_in_pdf, total_pages



class FolderZipper:
    def zip_folder(self,input_path,output_zip):
        """
        :param input_path:
        :param output_zip:
        :return:
        """
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"The input path does not exist: {input_path}")

        with zipfile.ZipFile(output_zip, 'w', zipfile.ZIP_DEFLATED) as zipf:
            if os.path.isfile(input_path):
                zipf.write(input_path, os.path.basename(input_path))

            elif os.path.isdir(input_path):
                for root, dirs, files in os.walk(input_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        zipf.write(file_path, os.path.relpath(file_path, input_path))
            else:
                raise ValueError

    def zip_folder_many(self, input_path, size, output_dir,remain_raw_name=False):
        """
        :param input_path:
        :param size:
        :param output_dir:
        :param remain_raw_name:
        :return:
        """
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        files = [os.path.join(input_path, f) for f in os.listdir(input_path)
                 if os.path.isfile(os.path.join(input_path, f)) and f.lower().endswith(('.png', '.jpg', '.jpeg'))]
        files=natsorted(files)

        if not remain_raw_name:
            for i in range(0, len(files), size):
                group = files[i:i + size]
                zip_filename = os.path.join(output_dir, f"{(i // size)}.zip")
                with zipfile.ZipFile(zip_filename, 'w') as zipf:
                    for file in group:
                        zipf.write(file, os.path.relpath(file, input_path))
        else:
            for file in files:
                zip_name_with_extention=os.path.basename(file)
                zip_name,_=os.path.splitext(zip_name_with_extention)
                zip_path = os.path.join(output_dir, f"{zip_name}.zip")
                self.zip_folder(file,zip_path)




class ExcelReaderOld:
    def set_file_path(self, file_path):
        """
        :param file_path:
        :return:
        """
        self.file_path = file_path

    def get_sheet_names(self):
        try:
            excel_file = pd.ExcelFile(self.file_path)
            sheet_names = excel_file.sheet_names
            print(f"文件中的 sheets: {sheet_names}")
            return sheet_names
        except Exception as e:
            print(f"无法获取 sheet 名称: {e}")
            return []

    def read_excel(self, sheet_name=None):
        """
        :param sheet_name:
        :return:
        """
        try:
            if self.file_path.endswith('.xls'):
                df = pd.read_excel(self.file_path, sheet_name=sheet_name, engine='xlrd')
            elif self.file_path.endswith('.xlsx'):
                df = pd.read_excel(self.file_path, sheet_name=sheet_name, engine='openpyxl')
            else:
                raise ValueError("文件类型不支持，请提供 .xls 或 .xlsx 文件。")

            print(f"成功读取 sheet: {sheet_name}")
            return df
        except Exception as e:
            print(f"读取文件时出错: {e}")
            return None

    def df_to_markdown(self, df, markdown_file_path):
        """
        :param df:
        :param markdown_file_path:
        :return:
        """
        markdown_data = df.to_markdown(index=False)
        with open(markdown_file_path, 'w', encoding='utf-8') as md_file:
            md_file.write(markdown_data)
        print(f"成功将数据转换为 Markdown 并保存到 {markdown_file_path}")


    def process_excel(self, excel_path, markdown_dir):
        """
        :param excel_path:
        :param markdown_dir:
        :return:
        """
        result = []
        self.set_file_path(excel_path)

        sheet_names = self.get_sheet_names()

        for sheet_name in sheet_names:
            df = self.read_excel(sheet_name=sheet_name)
            if df is not None and len(df)>0:
                markdown_file_path = f"{markdown_dir}/{sheet_name}.md"
                self.df_to_markdown(df, markdown_file_path)
                result.append({sheet_name: markdown_file_path})
            else:
                continue
        return result



class ExcelReader:
    def _check_path(self,excel_path,markdown_dir):
        """
        :param excel_path:
        :param markdown_dir:
        :return:
        """
        if not os.path.isfile(excel_path):
            raise FileNotFoundError(f"文件路径 {excel_path} 不存在，请提供有效的路径。")
        if not (excel_path.endswith('.xls') or excel_path.endswith('.xlsx')):
            raise ValueError("文件类型不支持，请提供 .xls 或 .xlsx 文件。")

        if not os.path.exists(markdown_dir):
            os.makedirs(markdown_dir)

    def get_sheet_names(self,file_path):
        """
        获取 Excel 文件中的 sheet 名称
        :return: sheet 名称列表
        """
        try:
            if file_path.endswith('.xlsx'):
                workbook = load_workbook(file_path)
                sheet_names = workbook.sheetnames
            else:
                workbook = xlrd.open_workbook(file_path)
                sheet_names = workbook.sheet_names()

            print(f"文件中的 sheets: {sheet_names}")
            return sheet_names
        except Exception as e:
            print(f"无法获取 sheet 名称: {e}")
            return []

    def read_excel_with_merged_cells(self,file_path, sheet_name=None):
        """
        读取包含合并单元格的 Excel sheet 并处理合并单元格
        :param sheet_name: sheet 名称
        :return: DataFrame 对象
        """
        try:
            if file_path.endswith('.xlsx'):
                # 处理 .xlsx 文件
                workbook = load_workbook(file_path)
                sheet = workbook[sheet_name]

                data = []
                for row in sheet.iter_rows(values_only=False):
                    row_data = []
                    for cell in row:
                        if cell.coordinate in sheet.merged_cells:
                            for merged_range in sheet.merged_cells.ranges:
                                if cell.coordinate in merged_range:
                                    cell_value = sheet[merged_range.coord.split(":")[0]].value
                                    row_data.append(cell_value)
                                    break
                        else:
                            row_data.append(cell.value)
                    data.append(row_data)
                df = pd.DataFrame(data)

            else:
                # 处理 .xls 文件
                workbook = xlrd.open_workbook(file_path)
                sheet = workbook.sheet_by_name(sheet_name)

                data = []
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    data.append(row)
                df = pd.DataFrame(data)

            if df.empty:
                print(f"Sheet '{sheet_name}' 为空，无数据。")
                return None

            print(f"成功读取 sheet: {sheet_name}")
            return df
        except Exception as e:
            print(f"读取文件时出错: {e}")
            return None

    def df_to_markdown(self, df, markdown_file_path):
        """
        将 DataFrame 转换为 Markdown 文件并保存
        :param df: DataFrame 对象
        :param markdown_file_path: Markdown 文件路径
        :return: None
        """
        try:
            markdown_data = df.to_markdown(index=False)
            with open(markdown_file_path, 'w', encoding='utf-8') as md_file:
                md_file.write(markdown_data)
            print(f"成功将数据转换为 Markdown 并保存到 {markdown_file_path}")
        except Exception as e:
            print(f"保存 Markdown 文件时出错: {e}")

    def process_excel(self, excel_path, markdown_dir):
        """
        处理 Excel 文件并将每个 sheet 转换为 Markdown 文件
        :param excel_path: Excel 文件路径
        :param markdown_dir: 保存 Markdown 文件的目录
        :return: 返回生成的 Markdown 文件路径列表
        """
        result = []
        try:
            self._check_path(excel_path,markdown_dir)

            sheet_names = self.get_sheet_names(excel_path)

            for sheet_name in sheet_names:
                df = self.read_excel_with_merged_cells(excel_path,sheet_name=sheet_name)
                if df is not None and not df.empty:
                    markdown_file_path = os.path.join(markdown_dir, f"{sheet_name}.md")
                    self.df_to_markdown(df, markdown_file_path)
                    result.append({sheet_name: markdown_file_path})
                else:
                    print(f"跳过 sheet: {sheet_name}，因为它为空或无法读取。")
                    continue
            return result
        except Exception as e:
            print(f"处理 Excel 文件时出错: {e}")
            return result


