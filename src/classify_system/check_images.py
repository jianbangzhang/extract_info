# !/usr/bin/env python
# -*-coding:utf-8 -*-

"""
# File       : check_images.py
# Time       ：2024/10/8 14:05
# Author     ：jianbang
# version    ：python 3.10
# company    : IFLYTEK Co.,Ltd.
# emil       : whdx072018@foxmail.com
# Description：
"""
import os
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import pdfplumber

class DocumentInspector(object):
    def check_images_tables_in_docx(self, docx_file):
        """
        :param docx_file: The path to the .docx file
        :return: List of tables found in the document
        """
        images_in_docx,tables_in_docx=[],[]
        doc = Document(docx_file)

        for i, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                images_in_docx.append(i)

        for i, table in enumerate(doc.tables):
            tables_in_docx.append(i)
        return images_in_docx+tables_in_docx,None

    def check_images_tables_in_pptx(self, pptx_file):
        """
        :param pptx_file: The path to the .pptx file
        :return: List of tables and images found in the presentation
        """
        images_in_pptx ,tables_in_pptx=[],[]
        prs = Presentation(pptx_file)

        total_pages=len(prs.slides)
        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # 13 indicates a picture
                    images_in_pptx.append(slide_index)
                elif shape.has_table:
                    tables_in_pptx.append(slide_index)
                else:
                    continue
        pages = sorted(list(set(images_in_pptx+tables_in_pptx)))
        return pages, total_pages


    def check_images_tables_in_pdf_old(self,pdf_file):
        """
        :param pdf_file: The path to the PDF file
        :return: List of tables found in the PDF
        """
        images_in_pdf, tables_in_pdf = [], []
        pdf = PdfReader(pdf_file)

        total_pages=len(pdf.pages)
        # Iterate through PDF pages
        for page_num, page in enumerate(pdf.pages):
            # Check for images in the PDF page
            if '/XObject' in page['/Resources']:
                x_object = page['/Resources']['/XObject']
                for obj in x_object:
                    if x_object[obj]['/Subtype'] == '/Image':
                        images_in_pdf.append(page_num)

            if '/Annots' in page:
                for annot_ref in page['/Annots']:
                    annot = annot_ref.get_object()  # Dereference the object
                    if annot.get('/Subtype') == '/Widget':
                        tables_in_pdf.append(page_num)

        pages=sorted(list(set(images_in_pdf + tables_in_pdf)))
        return pages,total_pages


    def check_images_tables_in_pdf(self, pdf_file):
        """
        :param pdf_file: The path to the PDF file
        :return: A tuple containing two lists - images found in the PDF, tables found in the PDF, and the total number of pages
        """
        images_in_pdf, tables_in_pdf = [], []

        with pdfplumber.open(pdf_file) as pdf:
            total_pages = len(pdf.pages)

            for page_num, page in enumerate(pdf.pages):
                width, height = page.width, page.height
                header_height = page.height / 10
                non_header_bbox = (0, header_height, width, height)

                cropped_page = page.within_bbox(non_header_bbox)

                if cropped_page.images:
                    images_in_pdf.append(page_num)

                if cropped_page.extract_tables():
                    tables_in_pdf.append(page_num)

        pages = sorted(list(set(images_in_pdf + tables_in_pdf)))
        return pages, total_pages





