# !/usr/bin/env python
# -*-coding:utf-8 -*-

"""
# File       : check_sim.py
# Time       ：2024/10/7 10:12
# Author     ：jianbang
# version    ：python 3.10
# company    : IFLYTEK Co.,Ltd.
# emil       : whdx072018@foxmail.com
# Description：
"""
import jieba.analyse
import hashlib
import pdfplumber
from docx import Document
from pptx import Presentation
import openpyxl
import os
import xlrd
from utils import read_files,get_path_info
from multiprocessing import Pool, cpu_count
from concurrent.futures import ThreadPoolExecutor, as_completed



class SimHash:
    def __init__(self, tokens):
        """
        :param tokens:
        """
        self.hashbits = 64
        self.simhash = self.simhash(tokens)

    def _string_hash(self, v):
        """
        :param v:
        :return:
        """
        if v == "":
            return 0
        x = int(hashlib.md5(v.encode('utf-8')).hexdigest(), 16)
        return x

    def simhash(self, tokens):
        """
        :param tokens:
        :return:
        """
        v = [0] * self.hashbits
        for t in tokens:
            token_hash = self._string_hash(t)
            for i in range(self.hashbits):
                bitmask = 1 << i
                if token_hash & bitmask:
                    v[i] += 1
                else:
                    v[i] -= 1
        fingerprint = 0
        for i in range(self.hashbits):
            if v[i] > 0:
                fingerprint |= 1 << i
        return fingerprint

    def hamming_distance(self, other):
        """
        :param other:
        :return:
        """
        x = self.simhash ^ other.simhash
        dist = 0
        while x:
            dist += 1
            x &= x - 1
        return dist

    def is_similar(self, other, threshold=3):
        return self.hamming_distance(other) <= threshold


import os
import csv
import pandas as pd
from pptx import Presentation
from docx import Document
import pdfplumber
import openpyxl
import xlrd
from bs4 import BeautifulSoup
from PIL import Image
import numpy as np


class FileContentReader(object):
    def read_txt(self, filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as txt_file:
                content = txt_file.read()
            return content
        except Exception as e:
            print(f"读取 TXT 文件 {filepath} 时出错: {e}")
            return ""

    def read_pdf(self, filepath):
        content = ''
        try:
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    content += page.extract_text() or ""
            return content
        except Exception as e:
            print(f"读取 PDF 文件 {filepath} 时出错: {e}")
            return ""

    def read_docx(self, filepath):
        try:
            doc = Document(filepath)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            print(f"读取 Word 文件 {filepath} 时出错: {e}")
            return ""

    def read_pptx(self, filepath):
        try:
            prs = Presentation(filepath)
            content = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + '\n'
            return content
        except Exception as e:
            print(f"读取 PPT 文件 {filepath} 时出错: {e}")
            return ""

    def read_xlsx(self, filepath):
        try:
            if filepath.endswith('.xls'):
                workbook = xlrd.open_workbook(filepath)
                sheet = workbook.sheet_by_index(0)
                content = ""
                for row in range(sheet.nrows):
                    content += "\t".join(map(str, sheet.row_values(row))) + "\n"
                return content
            else:
                workbook = openpyxl.load_workbook(filepath)
                sheet = workbook.active
                content = ""
                for row in sheet.iter_rows(values_only=True):
                    content += "\t".join(map(str, row)) + "\n"
                return content
        except Exception as e:
            print(f"读取 Excel 文件 {filepath} 时出错: {e}")
            return ""

    def read_csv(self, filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as csv_file:
                content = csv_file.read()
            return content
        except Exception as e:
            print(f"读取 CSV 文件 {filepath} 时出错: {e}")
            return ""

    def read_tsv(self, filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as tsv_file:
                content = tsv_file.read()
            return content
        except Exception as e:
            print(f"读取 TSV 文件 {filepath} 时出错: {e}")
            return ""

    def read_html(self, filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as html_file:
                soup = BeautifulSoup(html_file, 'html.parser')
                return soup.get_text()
        except Exception as e:
            print(f"读取 HTML 文件 {filepath} 时出错: {e}")
            return ""

    def read_image(self, filepath):
        print(f"图像文件读取接口尚未实现: {filepath}")
        return None  # Returns the image as a numpy array


    def read_audio(self, filepath):
        # Placeholder for audio processing; can use libraries like pydub or librosa
        print(f"音频文件读取接口尚未实现: {filepath}")
        return None

    def read_video(self, filepath):
        print(f"视频文件读取接口尚未实现: {filepath}")
        return None

    def read_file_content(self, filepath):
        _, file_extension = os.path.splitext(filepath)
        file_extension = file_extension.lower()
        try:
            if file_extension == '.pdf':
                return self.read_pdf(filepath)
            elif file_extension == '.docx':
                return self.read_docx(filepath)
            elif file_extension == '.pptx':
                return self.read_pptx(filepath)
            elif file_extension in ['.xlsx', '.xls']:
                return self.read_xlsx(filepath)
            elif file_extension in [".txt"]:
                return self.read_txt(filepath)
            elif file_extension == ".csv":
                return self.read_csv(filepath)
            elif file_extension == ".tsv":
                return self.read_tsv(filepath)
            elif file_extension == ".html":
                return self.read_html(filepath)
            elif file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.gif']:
                return self.read_image(filepath)  # This returns a numpy array
            elif file_extension in ['.mp3', '.wav', '.ogg']:
                return self.read_audio(filepath)
            elif file_extension in ['.mp4', '.avi', '.mov']:
                return self.read_video(filepath)
            else:
                print(f"跳过的文件类型: {file_extension}")
                return ""
        except Exception as e:
            print(f"读取文件 {filepath} 时出错: {e}")
            return ""



class DocumentProcessor:
    def __init__(self, config):
        """
        :param config:
        """
        self.filepaths = read_files(config.path_file)
        self.file_reader = FileContentReader()
        self.max_file_size = config.max_files_size
        self.unique_file_path = config.unique_file_path
        self.max_thread=config.max_thread
        self.skip_file_lst = []
        self.total_files = 0
        self.failed_lst = []
        self.set_runtime = config.set_runtime
        self.timeout = config.task_timeout


    def get_document_simhash(self, content):
        """
        :param content: Document content
        :return: SimHash object
        """
        keywords = jieba.analyse.extract_tags(content, topK=20, withWeight=False)
        return SimHash(keywords)

    def remove_duplicate_documents(self, documents_list):
        """
        :param documents: List of document contents
        :return: List of file paths for unique documents
        """
        unique_docs = []
        unique_filepaths = []
        doc_hashes = []

        for i, data_dict in enumerate(documents_list):
            path, doc = list(data_dict.items())[0]
            doc_simhash = self.get_document_simhash(doc)
            is_duplicate = False
            for doc_hash in doc_hashes:
                if doc_simhash.is_similar(doc_hash):
                    is_duplicate = True
                    break
            if not is_duplicate:
                unique_docs.append(doc)
                unique_filepaths.append(path)  # Save the file path of the unique document
                doc_hashes.append(doc_simhash)

        return unique_filepaths

    def process_single_file(self, filepath):
        """
        :param filepath: Path to the file
        :return: Content of the file or None if reading fails
        """
        doc_dict = {}
        print(f"Processing file: {filepath}")
        content = self.file_reader.read_file_content(filepath)
        if content is not None and len(content)>0:
            doc_dict[filepath]=content
            return doc_dict
        else:
            if filepath not in self.failed_lst:
                self.failed_lst.append(filepath)
            return None

    def write_to_txt(self, unique_filepaths):
        """
        :param unique_filepaths: List of file paths corresponding to unique documents
        """
        with open(self.unique_file_path, 'a', encoding='utf-8') as f:
            for i, filepath in enumerate(unique_filepaths):
                folder_path,filename,_=get_path_info(filepath)
                f.write(f"{filename}\t\t{filepath}" + "\n")
            print(f"File paths of unique documents written to: {self.unique_file_path}")

    def get_result(self):
        if self.total_files > 0:
            failure_rate = (len(self.failed_lst) / self.total_files) * 100
            print(f"总文件数: {self.total_files}")
            print(f"成功读取文件数: {self.total_files - len(self.failed_lst)}")
            print(f"未读取文件数: {len(self.failed_lst)},文件为：{self.failed_lst}")
            print(f"未读取的文件比例: {failure_rate:.2f}%")
        else:
            print("没有要处理的文件")


    def process_files(self):
        """
        Processes all files in the file paths list, removes duplicates, and writes the
        file paths of unique documents to the output file.
        """
        documents_list=[]
        self.total_files = len(self.filepaths)
        if self.total_files == 1:
            print("未发现重复文件")
            self.write_to_txt(self.filepaths)
            return

        if self.total_files > self.max_file_size:
            raise ValueError(
                f"Cannot process more than {self.max_file_size} files, reduce the number of files and try again!")

        for filepath in self.filepaths:
            documents = self.process_single_file(filepath)
            if documents:
                documents_list.append(documents)

        unique_filepaths = self.remove_duplicate_documents(documents_list)
        valid_filepaths=self.failed_lst+unique_filepaths
        self.write_to_txt(valid_filepaths)
        self.get_result()

    def process_files_with_accelerate(self):
        documents_list = []
        self.total_files = len(self.filepaths)
        if self.total_files==1:
            print("未发现重复文件")
            self.write_to_txt(self.filepaths)
            return

        if self.total_files > self.max_file_size:
            raise ValueError(f"Cannot process more than {self.max_file_size} files, reduce the number of files and try again!")

        # 根据文件数量和 CPU 动态决定线程数
        num_workers= min(min((cpu_count() + self.total_files) // 2, self.max_thread), self.total_files)
        print(f"Using {num_workers} threads for processing...")

        with ThreadPoolExecutor(max_workers=num_workers) as executor:
            futures = {executor.submit(self.process_single_file, filepath): filepath for filepath in self.filepaths}

            for future in as_completed(futures):
                filepath = futures[future]
                try:
                    if self.set_runtime:
                        documents = future.result(timeout=self.timeout)
                    else:
                        documents = future.result()

                    if documents:
                        documents_list.append(documents)
                except Exception as e:
                    print(f"Error processing file {filepath}: {e}")

        print(f"Successfully read {len(documents_list)} files, skip {len(self.failed_lst)} files.")
        unique_filepaths = self.remove_duplicate_documents(documents_list)
        valid_filepaths = list(set(self.failed_lst + unique_filepaths))
        self.write_to_txt(valid_filepaths)
        self.get_result()


